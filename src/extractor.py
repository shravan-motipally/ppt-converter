import pytesseract
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from io import BytesIO
from typing import IO


class PowerpointExtractor:
    def __init__(self, source_pptx: str | IO[bytes]):
        self.source_pptx = source_pptx

    def extract_images_and_text(self):
        # Create a presentation object
        prs = Presentation(self.source_pptx)

        # Create a dictionary to hold the text and images
        data = {}

        # Create a list to hold the text, images and notes
        text, images, notes = [], [], []

        # Create an index to keep track of the slide number
        idx = 0

        # Loop through each slide in the presentation
        for slide in prs.slides:
            idx += 1
            # Loop through each shape in the slide
            for shape in slide.shapes:
                # Check if the shape has text
                if shape.has_text_frame:
                    # Get the text from the shape
                    text.append(shape.text)

                # Check if the shape has an image
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Get the image from the shape
                    image = shape.image
                    image_bytes = image.blob
                    image_object = Image.open(BytesIO(image_bytes))
                    image_text = pytesseract.image_to_string(image_object)
                    image_to_save = { "image_text": image_text, "image_object": image_object }
                    images.append(image_to_save)

            # Check if the shape has notes
            if slide.has_notes_slide:
                # Get the notes from the shape
                notes.append(slide.notes_slide.notes_text_frame.text)

            # Add the text, images, and notes to the data dictionary
            data[f"slide_{idx}"] = {}
            data[f"slide_{idx}"]["text"] = text
            data[f"slide_{idx}"]["images"] = images
            data[f"slide_{idx}"]["notes"] = notes
            text, images, notes = [], [], []

        return data